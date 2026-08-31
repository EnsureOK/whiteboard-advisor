"""办公文件生成:工件 -> Excel / Word / PPT 真实文件(字节)。

- 检视矩阵 (review_matrix) -> .xlsx
- 文档类工件 (doc: 方案书/面谈提纲/跟进) -> .docx / .pptx
配色沿用工作台黑白精密体系:黑白灰为主,状态色只作语义。
"""

from __future__ import annotations

from io import BytesIO

# 状态 -> 填充色(浅底,打印友好)
_LEVEL_FILL = {"ok": "EAF4EE", "mid": "F7EFDF", "high": "F9EAE6"}
_LEVEL_FONT = {"ok": "14683E", "mid": "8A5406", "high": "A33823"}
_INK = "1A1A19"
_MUTED = "6E6E6B"
_LINE = "E7E7E4"


# ---------- Excel:保单检视矩阵 ----------

def review_matrix_xlsx(title: str, content: dict) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = "保单检视矩阵"

    thin = Side(style="thin", color=_LINE)
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    left = Alignment(horizontal="left", vertical="center", wrap_text=True)

    cols = content.get("cols", [])
    rows = content.get("rows", [])

    # 标题与摘要
    ws.cell(row=1, column=1, value=title).font = Font(size=14, bold=True, color=_INK)
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(2, len(cols) + 1))
    if content.get("summary"):
        c = ws.cell(row=2, column=1, value=content["summary"])
        c.font = Font(size=9, color=_MUTED)
        ws.merge_cells(start_row=2, start_column=1, end_row=2, end_column=max(2, len(cols) + 1))

    head_row = 4
    head_fill = PatternFill("solid", fgColor="F7F7F5")
    ws.cell(row=head_row, column=1, value="成员").font = Font(bold=True, size=10, color=_MUTED)
    ws.cell(row=head_row, column=1).fill = head_fill
    ws.cell(row=head_row, column=1).border = border
    for i, col in enumerate(cols, 2):
        c = ws.cell(row=head_row, column=i, value=col)
        c.font = Font(bold=True, size=10, color=_MUTED)
        c.fill = head_fill
        c.border = border
        c.alignment = center

    for ri, row in enumerate(rows, head_row + 1):
        c = ws.cell(row=ri, column=1, value=row.get("member", ""))
        c.font = Font(bold=True, size=10, color=_INK)
        c.border = border
        c.alignment = left
        for ci, col in enumerate(cols, 2):
            cell = (row.get("cells") or {}).get(col)
            wc = ws.cell(row=ri, column=ci)
            wc.border = border
            wc.alignment = center
            if not cell:
                wc.value = "—"
                wc.font = Font(size=10, color=_MUTED)
                continue
            detail = f"已有 {cell.get('current', 0) // 10000} 万 / 建议 {cell.get('need', 0) // 10000} 万"
            wc.value = f"{cell.get('text', '')}\n{detail}"
            level = cell.get("level", "")
            wc.font = Font(size=10, color=_LEVEL_FONT.get(level, _INK))
            if level in _LEVEL_FILL:
                wc.fill = PatternFill("solid", fgColor=_LEVEL_FILL[level])

    # 附注(未计入矩阵的保单)
    extra_row = head_row + len(rows) + 2
    extras = content.get("extras") or []
    if extras:
        ws.cell(row=extra_row, column=1, value="未计入矩阵的保单").font = Font(bold=True, size=10, color=_INK)
        for i, ex in enumerate(extras, 1):
            ws.cell(
                row=extra_row + i,
                column=1,
                value=f"{ex.get('line', '')} {ex.get('productName', '')} · 保额 {ex.get('amount', 0) // 10000} 万",
            ).font = Font(size=9, color=_MUTED)

    ws.column_dimensions["A"].width = 14
    for i in range(2, len(cols) + 2):
        ws.column_dimensions[get_column_letter(i)].width = 20
    for r in range(head_row + 1, head_row + 1 + len(rows)):
        ws.row_dimensions[r].height = 32

    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ---------- Word:文档类工件(方案书/面谈提纲/跟进) ----------

def doc_docx(title: str, content: dict) -> bytes:
    import docx
    from docx.shared import Pt, RGBColor

    d = docx.Document()

    h = d.add_heading(title, level=0)
    for run in h.runs:
        run.font.color.rgb = RGBColor.from_string(_INK)

    if content.get("summary"):
        p = d.add_paragraph(content["summary"])
        for run in p.runs:
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor.from_string(_MUTED)

    for sec in content.get("sections", []):
        if sec.get("heading"):
            hh = d.add_heading(sec["heading"], level=1)
            for run in hh.runs:
                run.font.color.rgb = RGBColor.from_string(_INK)
        body = sec.get("body", "")
        for para in body.split("\n"):
            para = para.strip()
            if not para:
                continue
            if para.startswith(("- ", "• ")):
                d.add_paragraph(para[2:], style="List Bullet")
            else:
                d.add_paragraph(para)

    footer = d.add_paragraph("由经纪人智能体工作台生成 · 内容以正式条款与核保结论为准")
    for run in footer.runs:
        run.font.size = Pt(8)
        run.font.color.rgb = RGBColor.from_string(_MUTED)

    buf = BytesIO()
    d.save(buf)
    return buf.getvalue()


# ---------- PPT:文档类工件的讲解版 ----------

def doc_pptx(title: str, content: dict) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor as PptxColor
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    ink = PptxColor.from_string(_INK)
    muted = PptxColor.from_string(_MUTED)

    # 封面
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.4))
    tf = tb.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = ink
    if content.get("summary"):
        sub = slide.shapes.add_textbox(Inches(0.92), Inches(4.1), Inches(11.5), Inches(0.9))
        stf = sub.text_frame
        stf.text = content["summary"]
        stf.paragraphs[0].font.size = Pt(14)
        stf.paragraphs[0].font.color.rgb = muted

    # 每个 section 一页
    for sec in content.get("sections", []):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        head = s.shapes.add_textbox(Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.9))
        htf = head.text_frame
        htf.text = sec.get("heading", "")
        htf.paragraphs[0].font.size = Pt(26)
        htf.paragraphs[0].font.bold = True
        htf.paragraphs[0].font.color.rgb = ink

        body = s.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(11.4), Inches(5.2))
        btf = body.text_frame
        btf.word_wrap = True
        lines = [ln.strip() for ln in (sec.get("body") or "").split("\n") if ln.strip()]
        for i, ln in enumerate(lines):
            p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            bullet = ln.startswith(("- ", "• "))
            p.text = ("• " + ln[2:]) if bullet else ln
            p.font.size = Pt(16)
            p.font.color.rgb = ink if not bullet else ink
            p.space_after = Pt(8)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------- 分发 ----------

EXPORTS: dict[str, dict[str, str]] = {
    # 工件类型 -> {fmt: 扩展名}
    "review_matrix": {"xlsx": "xlsx"},
    "plan_doc": {"docx": "docx", "pptx": "pptx"},
    "visit_outline": {"docx": "docx", "pptx": "pptx"},
    "followup_msg": {"docx": "docx"},
}

MEDIA_TYPES = {
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def export_artifact(artifact_type: str, title: str, content: dict, fmt: str) -> bytes:
    """按工件类型与格式生成文件字节;不支持的组合抛 ValueError。"""
    allowed = EXPORTS.get(artifact_type, {})
    if fmt not in allowed:
        raise ValueError(f"工件类型 {artifact_type} 不支持导出为 {fmt}")
    if artifact_type == "review_matrix":
        return review_matrix_xlsx(title, content)
    if fmt == "docx":
        return doc_docx(title, content)
    if fmt == "pptx":
        return doc_pptx(title, content)
    raise ValueError(f"unsupported fmt {fmt}")
